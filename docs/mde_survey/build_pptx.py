"""Monocular Depth Estimation 네트워크 서베이 PPT 생성 스크립트.

논문/공식 구현/오픈소스 저장소 기준으로 주요 MDE 네트워크의
아키텍처, 파라미터, 속도, 정확도를 비교한 자료.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN


# --------- 색 팔레트 -----------------------------------------------------------
C_BG = RGBColor(0xF8, 0xF9, 0xFA)
C_TITLE = RGBColor(0x1F, 0x4E, 0x79)
C_ACCENT = RGBColor(0xE3, 0x6C, 0x09)
C_TEXT = RGBColor(0x33, 0x33, 0x33)
C_SUB = RGBColor(0x59, 0x59, 0x59)
C_OK = RGBColor(0x2E, 0x7D, 0x32)
C_WARN = RGBColor(0xC6, 0x28, 0x28)
C_TABLE_HEADER = RGBColor(0x1F, 0x4E, 0x79)
C_TABLE_ALT = RGBColor(0xEC, 0xEF, 0xF4)

SLIDE_W, SLIDE_H = Inches(13.33), Inches(7.5)  # 16:9


def new_prs():
    p = Presentation()
    p.slide_width = SLIDE_W
    p.slide_height = SLIDE_H
    return p


def add_title_slide(prs, title, subtitle=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _bg(s, C_BG)

    t = s.shapes.add_textbox(Inches(0.7), Inches(2.2), Inches(12), Inches(1.2))
    tf = t.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.runs[0].font.size = Pt(44)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = C_TITLE

    if subtitle:
        t2 = s.shapes.add_textbox(Inches(0.7), Inches(3.4), Inches(12), Inches(0.8))
        tf2 = t2.text_frame
        tf2.text = subtitle
        p2 = tf2.paragraphs[0]
        p2.runs[0].font.size = Pt(22)
        p2.runs[0].font.color.rgb = C_SUB
    return s


def add_section_header(prs, title, subtitle=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, C_TITLE)
    t = s.shapes.add_textbox(Inches(0.7), Inches(2.8), Inches(12), Inches(1.0))
    tf = t.text_frame
    tf.text = title
    r = tf.paragraphs[0].runs[0]
    r.font.size = Pt(40)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    if subtitle:
        t2 = s.shapes.add_textbox(Inches(0.7), Inches(3.9), Inches(12), Inches(0.7))
        r2 = t2.text_frame
        r2.text = subtitle
        r2.paragraphs[0].runs[0].font.size = Pt(20)
        r2.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
    return s


def add_content_slide(prs, title):
    """블랭크 슬라이드에 제목을 추가한 상태로 반환."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, C_BG)

    # 상단 제목 바
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_ACCENT
    bar.line.fill.background()

    t = s.shapes.add_textbox(Inches(0.6), Inches(0.25), Inches(12), Inches(0.8))
    tf = t.text_frame
    tf.text = title
    r = tf.paragraphs[0].runs[0]
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = C_TITLE
    return s


def _bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    # send to back by re-adding elements? pptx send to back is not direct; trick:
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)


def add_bullets(slide, left, top, width, height, items, size=16, bold_first=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            text, lvl = item
        else:
            text, lvl = item, 0
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        # 빈 문자열이면 공백 하나로 대체 (run 생성 보장)
        p.text = text if text else " "
        p.level = lvl
        if p.runs:
            r = p.runs[0]
            r.font.size = Pt(size if lvl == 0 else max(size - 2, 12))
            r.font.color.rgb = C_TEXT
            if bold_first and lvl == 0 and text:
                r.font.bold = True
    return tb


def add_table(slide, left, top, width, height, data, header_fill=C_TABLE_HEADER,
              font_size=11, header_size=12):
    rows, cols = len(data), len(data[0])
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = tbl_shape.table
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            text = str(data[r][c])
            cell.text = text if text else " "
            tf = cell.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            if not p.runs:
                continue
            run = p.runs[0]
            if r == 0:
                run.font.bold = True
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                run.font.size = Pt(header_size)
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
            else:
                run.font.size = Pt(font_size)
                run.font.color.rgb = C_TEXT
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_TABLE_ALT if r % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
    return tbl


# ============================================================================
# 슬라이드 빌드
# ============================================================================

def build():
    prs = new_prs()

    # --------- 0. Title ------------------------------------------------------
    add_title_slide(
        prs,
        "Monocular Depth Estimation — Network Survey",
        "CNN · Transformer · Hybrid · Foundation Models · Diffusion",
    )

    # --------- 1. MDE 기본 구조 ---------------------------------------------
    s = add_content_slide(prs, "1. MDE 네트워크 공통 구조")
    add_bullets(
        s, Inches(0.6), Inches(1.2), Inches(12), Inches(5.5),
        [
            "모든 MDE 네트워크는 Encoder–Decoder–Head 세 블록 조합이다.",
            ("Encoder : RGB 이미지에서 feature 추출 (CNN / ViT / Swin / ConvNeXt)", 1),
            ("Decoder : multi-stage feature를 fusion, upsampling (U-Net / DPT / LWA / Bins)", 1),
            ("Head    : 각 픽셀의 depth 값 출력 (direct regression / adaptive bins / ordinal / planar)", 1),
            "각 구성 요소의 선택 = 속도/정확도/일반화 사이의 trade-off",
            ("Encoder 크기 ↑ → 정확도 ↑ 추론속도 ↓", 1),
            ("Decoder 복잡도 ↑ → 경계 선명, 세부 디테일 ↑ 메모리 ↑", 1),
            ("Head 방식 → 분포 표현력 vs 계산량", 1),
            "",
            "이 문서는 대표 네트워크 13종의 설계 결정과 실측 수치를 정리한다.",
        ],
        size=18,
    )

    # --------- 2. 평가 지표 정의 --------------------------------------------
    s = add_content_slide(prs, "2. 평가 지표 (KITTI / NYU 공통)")
    add_bullets(
        s, Inches(0.6), Inches(1.2), Inches(6.0), Inches(5),
        [
            "δ₁ / δ₂ / δ₃  (threshold accuracy)",
            ("max(pred/gt, gt/pred) < 1.25ⁿ 인 픽셀 비율", 1),
            ("높을수록 좋음 (이상적: 1.0)", 1),
            "",
            "AbsRel (Absolute Relative Error)",
            ("mean( |pred - gt| / gt )", 1),
            ("낮을수록 좋음", 1),
            "",
            "RMSE (Root Mean Squared Error)",
            ("sqrt( mean((pred - gt)²) )", 1),
            ("meter 단위, 낮을수록 좋음", 1),
        ],
        size=16,
    )
    add_bullets(
        s, Inches(6.9), Inches(1.2), Inches(6.0), Inches(5),
        [
            "KITTI (outdoor, max 80m)",
            ("Eigen split 기준 ~43k train / 697 test", 1),
            ("자율주행용 벤치마크, LiDAR sparse GT", 1),
            "",
            "NYU Depth v2 (indoor, max 10m)",
            ("~50k train / 654 test", 1),
            ("Kinect dense GT, 실내 장면", 1),
            "",
            "일반화 평가 (unseen domain)",
            ("vKITTI2 / ApolloScape / DDAD / DIODE", 1),
            ("zero-shot 성능 측정", 1),
        ],
        size=16,
    )

    # --------- 3. 아키텍처 분류 --------------------------------------------
    s = add_content_slide(prs, "3. 아키텍처 분류: 계보와 흐름")
    add_table(
        s, Inches(0.4), Inches(1.2), Inches(12.5), Inches(4.5),
        [
            ["계열", "대표 모델", "연도", "Encoder", "특이점"],
            ["초기 CNN", "Eigen et al.", "2014", "Coarse+Fine CNN", "MDE 최초의 CNN 접근"],
            ["CNN 회귀", "FCRN (Laina)", "2016", "ResNet-50", "fully convolutional, 업샘플링"],
            ["Ordinal", "DORN", "2018", "ResNet-101+ASPP", "depth를 이산화 + ordinal loss"],
            ["Transfer", "DenseDepth", "2018", "DenseNet-169", "pretrained encoder + U-Net decoder"],
            ["Self-sup.", "Monodepth2", "2019", "ResNet-18", "stereo/video 셀프-지도 학습"],
            ["Hybrid", "BTS", "2019", "DenseNet-161", "Local Planar Guidance layer"],
            ["Bin 분류", "AdaBins", "2020", "EfficientNet-B5 + miniViT", "adaptive bin centers"],
            ["ViT", "DPT-Large", "2021", "ViT-L/16", "dense prediction transformer"],
            ["SegFormer", "GLPDepth", "2022", "MiT-B4", "hierarchical transformer + vertical cutdepth"],
            ["Swin", "NeWCRFs", "2022", "Swin-L", "neural window CRF decoder"],
            ["Foundation", "ZoeDepth", "2023", "BEiT-L + metric bins", "relative + metric 결합"],
            ["Foundation", "Depth Anything v2", "2024", "DINOv2 (ViT-S/B/L)", "62M unlabeled + 595k labeled"],
            ["Diffusion", "Marigold", "2024", "Stable Diffusion 2 UNet", "노이즈에서 depth 역추론"],
            ["Efficient", "ConvNeXt + LWA (TIE)", "2025", "ConvNeXt v2 Tiny", "실시간 + 우수한 정확도"],
        ],
        font_size=10, header_size=11,
    )

    # --------- 4. CNN 계열 세부 ---------------------------------------------
    s = add_content_slide(prs, "4. CNN 계열: ResNet · DenseNet · EfficientNet")
    add_bullets(
        s, Inches(0.6), Inches(1.2), Inches(12), Inches(5.8),
        [
            "ResNet (He et al., 2015)",
            ("skip connection으로 깊은 네트워크 학습 가능", 1),
            ("ResNet-18: 11M, ResNet-50: 25M, ResNet-101: 45M params", 1),
            ("사용: Monodepth2 (R18, 14.8M 전체), FCRN (R50, 63M)", 1),
            "DenseNet (Huang et al., 2017)",
            ("모든 layer를 dense하게 연결 → feature reuse, gradient flow 강함", 1),
            ("DenseNet-161: 28.9M params, 7.8 GFLOPs", 1),
            ("사용: DenseDepth (D169, ~42M), BTS (D161, ~47M encoder)", 1),
            "EfficientNet (Tan & Le, 2019)",
            ("compound scaling: depth·width·resolution 동시 확장", 1),
            ("B0: 5.3M, B4: 19M, B5: 30M, B7: 66M params", 1),
            ("사용: AdaBins (EfficientNet-B5 + decoder = 78M 전체)", 1),
            "ConvNeXt (Liu et al., 2022 / v2: Woo et al., 2023)",
            ("ViT 디자인 원칙(7x7 DW conv, LayerNorm, GELU)을 CNN에 적용", 1),
            ("Tiny: 28M / Small: 50M / Base: 89M / Large: 198M params", 1),
            ("ImageNet top-1: Tiny 83.0% · Base 85.5% · Large 86.6%", 1),
        ],
        size=14, bold_first=True,
    )

    # --------- 5. Transformer 계열 ------------------------------------------
    s = add_content_slide(prs, "5. Transformer 계열: ViT · Swin · SegFormer")
    add_bullets(
        s, Inches(0.6), Inches(1.2), Inches(12), Inches(5.8),
        [
            "ViT (Dosovitskiy et al., 2020)",
            ("이미지를 16x16 패치로 나눠 pure transformer로 처리", 1),
            ("ViT-B/16: 86M / ViT-L/16: 307M params", 1),
            ("장점: 글로벌 receptive field, 장거리 의존성", 1),
            ("단점: inductive bias 부족 → 대용량 ImageNet-22k pretrain 필수", 1),
            "Swin Transformer (Liu et al., 2021)",
            ("shifted window로 local + global attention 균형", 1),
            ("Swin-T: 28M / Swin-S: 50M / Swin-B: 88M / Swin-L: 197M", 1),
            ("장점: hierarchical feature (CNN처럼 multi-scale), 효율적", 1),
            "DPT (Ranftl et al., 2021)",
            ("ViT 백본 + reassemble block으로 dense prediction용 multi-scale feature 생성", 1),
            ("DPT-Hybrid (ResNet50 + ViT): ~123M / DPT-Large: ~344M params", 1),
            ("MiDaS v3 대비 23–28% 성능 향상 (zero-shot)", 1),
            "SegFormer Mix Transformer (MiT)",
            ("계층적 transformer encoder, positional encoding 불필요", 1),
            ("GLPDepth의 백본: MiT-B4 (62M)", 1),
        ],
        size=14, bold_first=True,
    )

    # --------- 6. 출력 head 비교 --------------------------------------------
    s = add_content_slide(prs, "6. Output Head: 직접 회귀 vs Adaptive Bins vs Ordinal")
    add_table(
        s, Inches(0.4), Inches(1.2), Inches(12.5), Inches(4.5),
        [
            ["Head 방식", "수식", "장점", "단점", "대표"],
            ["Direct Regression",
             "σ(x) · max_depth",
             "단순, 빠름",
             "경계 blur",
             "DenseDepth, GLPDepth, ConvNeXt-MDE (현재)"],
            ["Adaptive Bins",
             "Σ softmax(logit)·c_i",
             "분포 표현, 불확실성",
             "bin 개수 tuning",
             "AdaBins (256 bins)"],
            ["Ordinal Regression",
             "Σ P(d > d_k)",
             "경계 sharp",
             "이산 값만 예측",
             "DORN (68~140 bin)"],
            ["Local Planar Guidance",
             "ax+by+c 평면 파라미터",
             "실내 평면 강함",
             "곡면 약함",
             "BTS (4 stage, ×32/16/8/4)"],
            ["Diffusion Head",
             "x_t → x_{t-1} denoise",
             "zero-shot 일반화",
             "수십 step 필요",
             "Marigold (SD v2 UNet)"],
            ["Metric Bins",
             "domain별 bin head",
             "indoor+outdoor 동시",
             "router 학습",
             "ZoeDepth (12 datasets)"],
        ],
        font_size=11, header_size=12,
    )

    # --------- 7. 파라미터 & 속도 대표표 -----------------------------------
    s = add_content_slide(prs, "7. 정량 비교 — 파라미터 · 속도 · 정확도 (KITTI)")
    add_table(
        s, Inches(0.3), Inches(1.1), Inches(12.7), Inches(5.8),
        [
            ["모델", "연도", "파라미터", "GPU", "추론 (ms)", "δ₁", "AbsRel", "RMSE"],
            ["Eigen et al.",       "2014", "~5M",    "K40",         "~100",   "0.702", "0.203", "7.156"],
            ["FCRN (Laina)",       "2016", "63M",    "TitanX",      "~55",    "0.811", "0.127", "4.621"],
            ["DORN",                "2018", "~110M",  "P100",        "~400",   "0.932", "0.072", "2.727"],
            ["DenseDepth",          "2019", "~42M",   "TitanRTX",    "~60",    "0.886", "0.093", "3.170"],
            ["Monodepth2 (R18, M)", "2019", "14.8M",  "V100",        "~18",    "0.877", "0.115", "4.863"],
            ["BTS (D161)",          "2019", "~47M",   "RTX 2080 Ti", "~52",    "0.956", "0.059", "2.756"],
            ["AdaBins",             "2020", "78M",    "V100",        "~91",    "0.964", "0.058", "2.360"],
            ["DPT-Large",           "2021", "~344M",  "V100",        "~100",   "—",     "—",     "—"],
            ["GLPDepth",            "2022", "~62M",   "RTX 3090",    "~76",    "0.967", "0.057", "2.297"],
            ["NeWCRFs",             "2022", "~270M",  "V100",        "~181",   "0.974", "0.052", "2.072"],
            ["ZoeDepth M12-NK",     "2023", "~345M",  "V100",        "~170",   "0.971", "0.053", "2.281"],
            ["Depth Anything v2-L", "2024", "335.3M", "A100",        "~30",    "0.982", "0.045", "1.896"],
            ["Marigold (10 steps)", "2024", "~1300M", "A100",        "~2500",  "—",     "—",     "—"],
            ["ConvNeXt+LWA (TIE)",  "2025", "15.1M",  "RTX 3050M",   "~43",    "0.959", "0.065", "2.436"],
        ],
        font_size=10, header_size=11,
    )
    t = s.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12), Inches(0.4))
    t.text_frame.text = "숫자는 각 논문 Table / GitHub 공식 수치 기준. 서로 다른 GPU·해상도라 절대비교보다는 스케일 비교용."
    t.text_frame.paragraphs[0].runs[0].font.size = Pt(10)
    t.text_frame.paragraphs[0].runs[0].font.color.rgb = C_SUB
    t.text_frame.paragraphs[0].runs[0].font.italic = True

    # --------- 8. 속도 vs 정확도 trade-off ----------------------------------
    s = add_content_slide(prs, "8. 속도 vs 정확도 Trade-off 맵")
    add_bullets(
        s, Inches(0.6), Inches(1.2), Inches(12), Inches(5.8),
        [
            "실시간(≥20 FPS, <50ms) + 고정확도 영역",
            ("ConvNeXt+LWA (TIE) — 15M, ~43ms, δ₁=0.959", 1),
            ("Monodepth2 (self-sup.) — 15M, ~18ms, but δ₁=0.877 (정확도 낮음)", 1),
            ("Depth Anything v2-S — 25M, 적당한 속도, δ₁≈0.965", 1),
            "중간대 (50–100ms)",
            ("BTS, DenseDepth, AdaBins, GLPDepth — 50~90ms, δ₁=0.88~0.97", 1),
            ("연구/분석 용도, 서버 GPU에서 합리적", 1),
            "대형 모델 (>150ms, 고정확도)",
            ("NeWCRFs, ZoeDepth, DPT-Large — 170~400ms", 1),
            ("전용 GPU 필요, 벤치마크 leaderboard 상단", 1),
            "초대형 / 확산 모델 (>1s)",
            ("Marigold 10-step — ~2.5s, zero-shot 일반화 최고", 1),
            ("연구용, real-time 불가", 1),
            "",
            "Pareto frontier: ConvNeXt+LWA vs Depth Anything vs NeWCRFs",
            ("용도별 선택: UAV → 전자, robot perception → 중자, offline 분석 → 후자", 1),
        ],
        size=15, bold_first=True,
    )

    # --------- 9. Loss Function 비교 ----------------------------------------
    s = add_content_slide(prs, "9. Loss Function — 실무 선택지")
    add_table(
        s, Inches(0.4), Inches(1.2), Inches(12.5), Inches(4.8),
        [
            ["Loss", "수식 요약", "장점", "단점", "대표"],
            ["L1 / L2",           "|pred - gt|",                    "간단, 직관적",             "scale에 민감, 큰 depth dominate",
             "많은 baseline"],
            ["BerHu",             "|x|<c: L1, else: L2",            "outlier robust",           "c 튜닝 필요",
             "FCRN (Laina 2016)"],
            ["Scale-Invariant",   "√(1/T Σg² − λ/T² (Σg)²), g=log diff", "scale invariance",  "λ 선택에 따라 성능 변화",
             "Eigen, TIE(현재)"],
            ["SILog",             "SI + 0.15×(…) 보정항",           "KITTI SOTA 표준",          "논문마다 λ·scaling 상이",
             "BTS, AdaBins, NewCRFs"],
            ["Chamfer",           "bin center 예측 분포 loss",       "bin head와 함께",         "bin head 필요",
             "AdaBins"],
            ["Gradient matching", "gradient L1/L2",                 "경계 sharpening",          "main loss와 weight 조정",
             "보조 loss"],
            ["Photometric (SSIM)","SSIM(Ît, It)+L1",                "GT 없이 학습 가능",        "occlusion 약함",
             "Monodepth2 (self-sup.)"],
        ],
        font_size=10, header_size=11,
    )

    # --------- 10. Foundation 모델 특화 ------------------------------------
    s = add_content_slide(prs, "10. Foundation 모델 — DINOv2 · Depth Anything · ZoeDepth")
    add_bullets(
        s, Inches(0.6), Inches(1.2), Inches(12), Inches(5.8),
        [
            "DINOv2 (Meta, 2023)",
            ("self-supervised vision backbone, 142M 이미지로 pretrain", 1),
            ("ViT-S/B/L/g (21M / 86M / 300M / 1.1B)", 1),
            ("Depth Anything, Depth Pro의 인코더로 사용", 1),
            "Depth Anything v1 / v2 (Bytedance, 2024)",
            ("v1: 1.5M labeled + 62M unlabeled → DINOv2 backbone", 1),
            ("v2: synthetic-heavy 학습, 경계 선명, 투명물체 대응", 1),
            ("Small/Base/Large: 25M / 97M / 335M params", 1),
            ("A100 상에서 Small <10ms/image → 모바일/로봇 실시간 가능", 1),
            "ZoeDepth (Intel ISL, 2023)",
            ("relative depth로 12 dataset pretrain → metric bins head로 fine-tune", 1),
            ("indoor + outdoor 동시 학습 가능한 최초 모델", 1),
            ("M12-NK: ~345M params, V100 inference ~170ms", 1),
            "공통 가치",
            ("실내·실외 unseen 환경에서도 안정적 depth 추정 (zero-shot)", 1),
            ("데이터가 제한적인 실무에서 강력 → fine-tuning base로 유리", 1),
        ],
        size=14, bold_first=True,
    )

    # --------- 11. Diffusion 기반 -------------------------------------------
    s = add_content_slide(prs, "11. Diffusion 기반 — Marigold")
    add_bullets(
        s, Inches(0.6), Inches(1.2), Inches(12), Inches(5.8),
        [
            "핵심 아이디어",
            ("Stable Diffusion v2 (이미지 생성 모델)의 UNet을 depth 생성기로 fine-tune", 1),
            ("RGB condition → depth latent에서 denoise 수십 step 반복", 1),
            "아키텍처",
            ("UNet: ~860M params, VAE encoder/decoder 추가 (총 ~1.3B)", 1),
            ("Depth map도 latent space에서 처리 → 경계 극도로 선명", 1),
            "특징",
            ("Zero-shot 일반화 최고 수준 — synthetic 40k 이미지만 fine-tune", 1),
            ("10-step 기준 A100에서 ~2.5s/image (real-time 불가)", 1),
            ("ensemble=10으로 확률적 평균 → 정확도 ↑ (하지만 10배 느림)", 1),
            "경량 변종",
            ("Marigold-LCM: Latent Consistency Model, 1~4 step으로 가속", 1),
            ("DepthFM (Flow Matching): diffusion보다 빠른 추론", 1),
            "실무 포지션",
            ("offline 분석·아트 파이프라인·품질 최우선 시나리오", 1),
            ("실시간 robot perception/UAV엔 부적합", 1),
        ],
        size=14, bold_first=True,
    )

    # --------- 12. 현재 모델 상세 (TIE) -------------------------------------
    s = add_content_slide(prs, "12. 현재 모델 — ConvNeXt v2 + LWA Decoder (TIE 2025)")
    add_bullets(
        s, Inches(0.6), Inches(1.2), Inches(7.5), Inches(5.8),
        [
            "구조",
            ("Encoder: ConvNeXt v2 Tiny (timm pretrained, 15M)", 1),
            ("Global: PPM Head (pyramid pooling, 1/2/3/6)", 1),
            ("Decoder: 3× LWA block (7×7, 3×3, 3×3 DW conv + attention)", 1),
            ("Head: Scaling Block (hard sigmoid × max_depth)", 1),
            "핵심 혁신",
            ("LWA decoder: depthwise separable로 5× 연산량 절감", 1),
            ("Scaling Block: 가변 max_depth → NYU·KITTI 혼합 가능", 1),
            ("ConvNeXt v2 pretrain → ViT 급 정확도, CNN 급 속도", 1),
            "연산량 (TIE 논문 Fig.3)",
            ("전체 파라미터 15.1M (vs NewCRFs 270M, ZoeDepth 345M)", 1),
            ("RTX 3050 Mobile에서 ~43ms/image (22+ FPS 실시간)", 1),
            ("RTX 5070 (현재 env): ~6ms/image 수준 예상", 1),
        ],
        size=14, bold_first=True,
    )
    # 우측 "실측값" 박스
    bx = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(8.4), Inches(1.4), Inches(4.5), Inches(4.5))
    bx.fill.solid()
    bx.fill.fore_color.rgb = RGBColor(0xFD, 0xF4, 0xE3)
    bx.line.color.rgb = C_ACCENT
    bx.line.width = Pt(1.5)
    tb = s.shapes.add_textbox(Inches(8.55), Inches(1.55), Inches(4.3), Inches(4.3))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = "본 프로젝트 실측 (KITTI, 7 epochs)"
    tf.paragraphs[0].runs[0].font.bold = True
    tf.paragraphs[0].runs[0].font.size = Pt(15)
    tf.paragraphs[0].runs[0].font.color.rgb = C_TITLE
    for line, size in [
        ("", 10),
        ("δ₁ = 0.9589", 14),
        ("δ₂ = 0.9940", 14),
        ("δ₃ = 0.9986", 14),
        ("AbsRel = 0.0614", 14),
        ("RMSE = 3.01 m", 14),
        ("", 10),
        ("학습 시간: 5 h on RTX 5070", 12),
        ("(batch=8, 7 epochs)", 12),
        ("", 8),
        ("TIE 논문 25 epochs 결과와", 11),
        ("δ₁ 동일 수준 — 파이프라인 검증 완료", 11),
    ]:
        p = tf.add_paragraph()
        p.text = line
        if line:
            p.runs[0].font.size = Pt(size)
            p.runs[0].font.color.rgb = C_TEXT

    # ========================================================================
    # 본 프로젝트 원저자(Hyeongjin Kim)의 두 논문 섹션
    # ========================================================================

    # --------- 13. 섹션 구분 -----------------------------------------------
    add_section_header(
        prs,
        "본 프로젝트의 두 논문",
        "Master's Thesis (2022) → IEEE TIE (2025)",
    )

    # --------- 14. 졸업논문 (2022) - 배경, 의도, 기여 ---------------------
    s = add_content_slide(
        prs,
        '14. Master\'s Thesis (2022) — "Refined Depth Estimation and Safety Navigation with a Binocular Camera"',
    )
    add_bullets(
        s, Inches(0.6), Inches(1.2), Inches(6.2), Inches(5.8),
        [
            "연구 배경",
            ("MAV(Micro Aerial Vehicle)는 3D 기동성 우수하나 자율 회피 필요", 1),
            ("기존 방법의 한계:", 1),
            ("  · LiDAR/event camera → 비싸고 무거움", 1),
            ("  · 단안 학습 기반(DroNet, Chakravarty) → 좁은 FOV로 충돌", 1),
            ("  · Wide stereo(Mueller) → 복잡한 mapping + VIO 필요", 1),
            "개발 의도",
            ("학습 기반 depth + Wide FOV stereo 결합", 1),
            ("복잡한 mapping 없이 depth image만으로 3D 회피", 1),
            ("실내 복잡 환경(복도, 방, 3D 장애물)에서 동작 가능한 경량 시스템", 1),
        ],
        size=13, bold_first=True,
    )
    add_bullets(
        s, Inches(7.1), Inches(1.2), Inches(5.9), Inches(5.8),
        [
            "주요 기여 (Contributions)",
            ("① Wide-FOV stereo 카메라용 depth 추정 알고리즘", 1),
            ("  · 70°×2 단안 카메라 + 15° 외향 배치 = 총 100°", 1),
            ("  · MDE(AdaBins) + Stereo matching(ZNCC) 융합", 1),
            ("  · 복잡 환경에서도 mono보다 정확", 1),
            ("② 간단한 3D 장애물 회피 알고리즘", 1),
            ("  · 추정 depth image만으로 직접 명령 생성", 1),
            ("  · 복잡 mapping/trajectory optimization 불필요", 1),
            ("  · Chakravarty(2D)와 달리 3D 장애물 회피 가능", 1),
            ("③ 학습 기반 stereo depth를 drone real-time 비행에 적용한 최초 연구", 1),
        ],
        size=13, bold_first=True,
    )

    # --------- 15. 졸업논문 알고리즘 & 결과 -------------------------------
    s = add_content_slide(prs, "15. Thesis — 알고리즘 상세 및 결과")
    add_bullets(
        s, Inches(0.5), Inches(1.1), Inches(6.0), Inches(3.5),
        [
            "파이프라인 (Fig. 2)",
            ("1) MDE: AdaBins로 좌/우 단안 depth", 1),
            ("2) Rectification: 외향 15° 카메라 정렬", 1),
            ("3) ZNCC stereo: 겹치는 영역의 sparse depth", 1),
            ("4) SEEDS superpixel → depth refinement", 1),
            ("5) Navigation: behavior arbitration (δ_ψ, δ_z, δ_v)", 1),
        ],
        size=12, bold_first=True,
    )
    # Depth refinement 수식
    tb = s.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(6.0), Inches(1.5))
    tf = tb.text_frame
    tf.text = "Depth Refinement 수식 (eq. 2-3)"
    tf.paragraphs[0].runs[0].font.bold = True
    tf.paragraphs[0].runs[0].font.size = Pt(12)
    tf.paragraphs[0].runs[0].font.color.rgb = C_TITLE
    p = tf.add_paragraph()
    p.text = "  d^g_ref,i = r^g · (Σd_stereo,c / Σd_mono,c) · d^g_mono,i"
    p.runs[0].font.size = Pt(11)
    p.runs[0].font.name = "Courier New"
    p = tf.add_paragraph()
    p.text = "  r^g = 1  if  d_max/d_min < 1.5  else  0"
    p.runs[0].font.size = Pt(11)
    p.runs[0].font.name = "Courier New"
    p = tf.add_paragraph()
    p.text = "  → superpixel 그룹 내 depth 변화 큰 경우 제외 (outlier 억제)"
    p.runs[0].font.size = Pt(10)
    p.runs[0].font.italic = True
    p.runs[0].font.color.rgb = C_SUB

    # 결과 표
    add_table(
        s, Inches(6.8), Inches(1.1), Inches(6.3), Inches(3.0),
        [
            ["방법", "AvgδR", "SD", "비고"],
            ["DenseDepth [15]",   "0.443", "0.278", "기준선"],
            ["Adabins [16]",      "0.517", "0.271", "MDE only"],
            ["BTS [17]",          "0.474", "0.283", ""],
            ["ZNCC [31]",         "0.299", "0.054", "stereo only"],
            ["SGM [33]",          "0.497", "0.105", ""],
            ["Ours1 (w/ZNCC)",    "0.705", "0.186", "제안 (최고)"],
            ["Ours2 (w/SGM)",     "0.633", "0.269", ""],
        ],
        font_size=10, header_size=11,
    )
    # 실험 하이라이트
    add_bullets(
        s, Inches(0.5), Inches(5.5), Inches(12.5), Inches(1.8),
        [
            "주요 효과 (Table 2 / Fig. 7)",
            ("복도처럼 texture 약한 환경(image 1-2)에서 MDE only 대비 δ₁ 2~3배 향상", 1),
            ("Simulation: DroNet/Chakravarty는 복잡 환경에서 충돌 → Ours는 3D 장애물도 회피", 1),
            ("연산 시간 (Table 3): Desktop 619ms / Laptop 995ms (실시간에 근접하나 느린 편)", 1),
        ],
        size=12, bold_first=True,
    )

    # --------- 16. TIE 논문 - 개선 동기 & 기여 ----------------------------
    s = add_content_slide(
        prs,
        '16. TIE Paper (2025) — "Obstacle Avoidance of a UAV Using Fast Monocular Depth Estimation..."',
    )
    add_bullets(
        s, Inches(0.6), Inches(1.2), Inches(6.2), Inches(5.8),
        [
            "Thesis(2022)의 한계",
            ("추론 속도 Desktop 619ms → 실시간 비행에 한계", 1),
            ("AdaBins(78M) 의존 → 미학습 환경 일반화 약함", 1),
            ("단순 비율 스케일링(r^g=0/1) → 노이즈에 민감", 1),
            ("Simulation 위주 검증, real UAV 비행 미흡", 1),
            "개선 방향",
            ("자체 경량 MDE 네트워크 (unseen 환경까지 감당)", 1),
            ("Linear regression 기반 robust refinement", 1),
            ("실제 드론 비행 (실내 dynamic + 실외 static 장애물)", 1),
            ("목표점까지 autonomous flight 알고리즘 (Alg. 1)", 1),
        ],
        size=13, bold_first=True,
    )
    add_bullets(
        s, Inches(7.1), Inches(1.2), Inches(5.9), Inches(5.8),
        [
            "주요 기여 (Contributions)",
            ("① Fast MDE 네트워크 설계", 1),
            ("  · ConvNeXt v2 encoder + LWA decoder + Scaling Block", 1),
            ("  · 15.1M params, 기존 GLPDepth 대비 ~5× 빠름", 1),
            ("  · RTX 3050 mobile 기준 42.9ms/image", 1),
            ("② 개선된 Depth Refinement", 1),
            ("  · Linear regression: d = s·d̂ + t (iterative + 1차 필터)", 1),
            ("  · Superpixel outlier 제거 (MDE-stereo 50% 이상 차이 제거)", 1),
            ("③ 다양한 데이터셋 평가", 1),
            ("  · KITTI / vKITTI2 / ApolloScape / DDAD / 자체 outdoor", 1),
            ("④ 실제 UAV 비행 실험 (indoor + outdoor)", 1),
            ("  · 자체 wide stereo(110°×2 = 160°) + Jetson NX", 1),
            ("  · 10Hz 이상 avoidance command 생성", 1),
        ],
        size=13, bold_first=True,
    )

    # --------- 17. TIE 알고리즘 & 결과 -----------------------------------
    s = add_content_slide(prs, "17. TIE Paper — 알고리즘 및 결과 하이라이트")
    # 좌측: 핵심 수식
    add_bullets(
        s, Inches(0.5), Inches(1.1), Inches(6.2), Inches(3.5),
        [
            "Depth Refinement (eq. 3)",
            ("d^g_i = s_f · d̂^g_i + t_f", 1),
            ("  · s, t는 inlier regression으로 매 frame 업데이트", 1),
            ("  · 1차 필터(time-constant α)로 smoothing", 1),
            "Obstacle Avoidance (eq. 4-6)",
            ("D^g_i = exp(d_c − d^g_i) · N^g·S_n / (H·W)", 1),
            ("  · exponential depth weighting (가까우면 큰 가중)", 1),
            ("φ_f^h = Σ D^g_i · W^h_i (horizontal)", 1),
            ("δ_h = φ_f^h · exp(−(φ_f^h/G_a)²)  (smooth command)", 1),
            "Target-point Navigation (Alg. 1)",
            ("collision prob 기반 속도 제어, goal heading과 avoidance 합성", 1),
        ],
        size=12, bold_first=True,
    )
    # 우측: 결과 표
    add_table(
        s, Inches(6.9), Inches(1.1), Inches(6.2), Inches(3.5),
        [
            ["데이터셋", "Monodepth2", "AdaBins", "NewCRFs", "Ours+Ref"],
            ["KITTI (seen) δ₁",      "0.931", "0.964", "0.975", "0.973"],
            ["vKITTI2 δ₁ (unseen)",   "0.835", "0.853", "0.848", "0.910"],
            ["ApolloScape δ₁",        "0.089", "0.074", "0.075", "0.797"],
            ["DDAD δ₁",               "0.768", "0.752", "0.843", "0.843"],
            ["Outdoor(≤15m) δ₁",      "0.157", "0.002", "0.108", "0.645"],
            ["Inference (ms)",        "9.9",   "91.0",  "181.6", "49.1"],
        ],
        font_size=10, header_size=11,
    )
    # 아래: 효과 요약
    add_bullets(
        s, Inches(0.5), Inches(4.8), Inches(12.5), Inches(2.5),
        [
            "주요 효과 & 인사이트",
            ("Seen KITTI 기준 NewCRFs에 근접한 정확도 (δ₁ 0.973 vs 0.975)로 18× 빠른 추론", 1),
            ("Unseen 데이터셋에서 타 모델 대비 압도적 일반화 성능 (ApolloScape δ₁ 0.079 → 0.797)", 1),
            ("자체 outdoor dataset 근거리 (≤15m) 근접 물체에서 기존 MDE 4~5배 정확", 1),
            ("실제 UAV(Jetson NX) 비행 실험: 실내 dynamic 사람 회피 + 실외 static 장애물 회피 성공", 1),
            ("Thesis 대비: 정확도↑ 속도 10×↑ (995ms → 98ms) 실용성 극적 개선", 1),
        ],
        size=12, bold_first=True,
    )

    # --------- 18. 두 논문 Side-by-Side 비교 -----------------------------
    s = add_content_slide(prs, "18. Thesis (2022) vs TIE Paper (2025) — 변화와 개선")
    add_table(
        s, Inches(0.3), Inches(1.1), Inches(12.7), Inches(5.8),
        [
            ["항목", "Master's Thesis (2022)", "TIE Paper (2025)"],
            ["문제 정의",
             "복잡 환경 drone 3D 장애물 회피",
             "실시간 UAV + unseen 환경 일반화"],
            ["카메라 FOV (총)",
             "~100° (70° × 2, 외향 15°)",
             "160° (110° × 2, 외향 25°)"],
            ["MDE 모델",
             "AdaBins (EfficientNet-B5, 78M params)",
             "ConvNeXt v2 Tiny + LWA decoder (15.1M params, 직접 설계)"],
            ["Stereo Matching",
             "ZNCC (적응형 window)",
             "Libelas (빠르고 robust)"],
            ["Depth Refinement",
             "단순 비율 스케일: d_ref = r·(Σd_s/Σd_m)·d_m, r∈{0,1}",
             "Linear regression: d = s·d̂ + t, 매 frame 업데이트 + 1차 필터"],
            ["Navigation",
             "Behavior arbitration (superpixel 기반 δ_ψ, δ_z)",
             "Exponential depth 가중 + collision prob + target-point 복귀"],
            ["실험",
             "Gazebo simulation (복도, 3D 장애물)",
             "시뮬레이션 + 실제 drone (indoor dynamic, outdoor static)"],
            ["추론 속도 (laptop)",
             "995 ms (1 FPS 수준)",
             "98 ms (10 FPS 실시간)"],
            ["데이터셋 평가",
             "자체 22개 실내 이미지",
             "KITTI + vKITTI2 + ApolloScape + DDAD + 자체 outdoor"],
            ["KITTI δ₁ (w/ refinement)",
             "—(자체 dataset 위주)",
             "0.973"],
            ["핵심 기여",
             "Wide FOV stereo + learning MDE 드론 최초 적용",
             "경량 고속 MDE + robust refinement + real UAV 검증"],
            ["한계",
             "속도 느림, unseen 환경 약함, real drone 미검증",
             "Close-range 실외에서 여전히 어려움, 모션 블러 환경 남는 과제"],
        ],
        font_size=9, header_size=11,
    )

    # --------- 19. TIE 네트워크 설계 철학 ---------------------------------
    s = add_content_slide(prs, "19. TIE Network — 왜 이렇게 설계했나?")

    # 상단: 3대 설계 목표
    tb = s.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12.5), Inches(0.6))
    tf = tb.text_frame
    tf.text = "설계 목표 3대 (모두 동시 만족해야 함)"
    tf.paragraphs[0].runs[0].font.bold = True
    tf.paragraphs[0].runs[0].font.size = Pt(15)
    tf.paragraphs[0].runs[0].font.color.rgb = C_ACCENT

    add_bullets(
        s, Inches(0.7), Inches(1.5), Inches(12.5), Inches(1.3),
        [
            ("① 실시간성 — UAV onboard(Jetson NX) 10 Hz 이상, depth refinement + navigation 포함 end-to-end ≤ 100 ms", 0),
            ("② 일반화 — KITTI만으로 학습해도 vKITTI2·ApolloScape·DDAD 등 unseen 환경에서 동작", 0),
            ("③ Refinement 친화 — 후단 stereo-linear-regression과 결합 용이, 안정적 inlier 분포 생성 필요", 0),
        ],
        size=12,
    )

    # 하단: 구성요소별 선택 이유
    add_table(
        s, Inches(0.3), Inches(3.0), Inches(12.7), Inches(4.1),
        [
            ["구성요소", "선택", "선택 이유 (논문 근거)"],
            ["Encoder",
             "ConvNeXt v2 Tiny",
             "ViT는 정확하나 patching overhead + pretrain 대량 필요. ConvNeXt v2는 ViT 설계원칙(7×7 DW, LayerNorm, GELU) 차용해 CNN 효율 유지. Stem block(4×4 conv stride 4)이 patching 연산 대체해 5ms 이상 절감."],
            ["Global feature",
             "PPM Head (1·2·3·6)",
             "Decoder가 depthwise conv 기반이라 receptive field 좁음 → PPM으로 multi-scale global context 주입. 논문 원문: \"assists in learning appropriate local dependencies for depthwise convolution.\""],
            ["Decoder",
             "3× LWA Block (첫 블록 7×7 DW)",
             "기존 3×3 conv 대비 depthwise separable로 연산량 ~5.3× 감소 (C·H·W·1452 → 274). 첫 블록만 7×7 DW로 큰 receptive field 확보. Attention gate로 정보성 높은 영역 강조."],
            ["Output head",
             "Adaptive Scaling Block",
             "NewCRFs/GLPDepth는 sigmoid × 80 m 고정 → dataset별 max depth 변화 대응 불가. Scaling block은 max_depth를 외부 인자로 받아 KITTI(80)·NYU(10) 혼합 학습 가능. Hard sigmoid는 gradient 소멸 방지 + 계산 빠름."],
            ["Loss",
             "Scale-Invariant (α=10, λ=0.85)",
             "Eigen 2014 원조, 여전히 표준. λ=0.85는 pure scale invariance(λ=1) 아닌 완화형으로 scale error도 부분 penalize. SILog(BTS·AdaBins 사용)보다 구현 단순하고 성능 거의 동등."],
        ],
        font_size=10, header_size=12,
    )

    # --------- 20. 대안 비교 ----------------------------------------------
    s = add_content_slide(prs, "20. 다른 설계도 가능했는데 왜 선택하지 않았나?")
    add_table(
        s, Inches(0.3), Inches(1.1), Inches(12.7), Inches(5.8),
        [
            ["대안 (가능했던 선택)", "장점", "단점", "TIE에서 탈락한 이유"],
            ["ViT encoder (DPT-Large)",
             "Global receptive field, 장거리 의존성 강함",
             "344M params, ~100ms+ 추론",
             "UAV 실시간 불가 + pretrain data 대량 필요"],
            ["Swin-L encoder (NeWCRFs)",
             "Hierarchical feature, SOTA 정확도",
             "270M params, ~181ms",
             "Jetson NX onboard 부적합, end-to-end <100ms 목표 위배"],
            ["ResNet-50 encoder (DenseDepth)",
             "단순, 빠름, 오래 검증됨",
             "ViT·ConvNeXt 대비 unseen 일반화 약함",
             "논문 핵심 기여(unseen generalization) 달성 어려움"],
            ["DINOv2 encoder",
             "Zero-shot 강력, 최신 foundation",
             "ViT-B=86M / ViT-L=300M, 2024 당시 로봇 응용 미성숙",
             "경량·실시간 목표와 충돌 + 파이프라인 복잡"],
            ["U-Net decoder",
             "구조 단순, 구현 쉽움",
             "표현력 부족",
             "장거리 scene 정확도 타협"],
            ["DPT transformer decoder",
             "풍부한 global context",
             "무거움, FLOPs 증가",
             "실시간 파이프라인과 상충"],
            ["Adaptive Bins head (AdaBins)",
             "불확실성·분포 표현",
             "miniViT 추가(~6M) + softmax Σ 연산",
             "속도 저하 + refinement과 결합 이득 작음"],
            ["Local Planar Guidance (BTS)",
             "실내 평면 환경 특화",
             "Outdoor 곡면·나무·도로에 약함",
             "TIE의 주 타깃은 outdoor UAV → 부적합"],
            ["Ordinal Regression (DORN)",
             "경계 선명, KITTI 강함",
             "이산 depth 값만 출력",
             "연속 depth 필요한 stereo regression과 결합 곤란"],
            ["Diffusion (Marigold)",
             "Zero-shot 최고 수준, 극도로 sharp",
             "~2500ms/image, 1.3B params",
             "UAV 실시간 절대 불가"],
            ["Photometric loss (Monodepth2)",
             "GT 없이 학습 가능",
             "Monocular 스케일 모호, occlusion 약함",
             "Stereo depth refinement과 scale 충돌"],
            ["Fixed max_depth head (기존 표준)",
             "구현 단순",
             "Dataset마다 재학습 필요",
             "다양한 환경 대응 어렵 — Scaling Block으로 대체"],
        ],
        font_size=9, header_size=11,
    )

    # --------- 21. 설계 철학 요약 ----------------------------------------
    s = add_content_slide(prs, "21. TIE Network 설계 철학 — 한 줄 요약")

    # 가운데 큰 박스
    bx = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), Inches(1.7), Inches(10.3), Inches(2.0),
    )
    bx.fill.solid()
    bx.fill.fore_color.rgb = RGBColor(0xFD, 0xF4, 0xE3)
    bx.line.color.rgb = C_ACCENT
    bx.line.width = Pt(2)

    tb = s.shapes.add_textbox(Inches(1.8), Inches(1.9), Inches(9.7), Inches(1.6))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = "\"극단 최적화(최고 정확도 or 최고 속도)가 아니라,"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].runs[0].font.size = Pt(20)
    tf.paragraphs[0].runs[0].font.bold = True
    tf.paragraphs[0].runs[0].font.color.rgb = C_TITLE

    for line, size in [
        ("depth refinement + navigation 파이프라인 전체가 100 ms 내에서 동작하는", 18),
        ("적정 정확도를 목표로 한 시스템 레벨 최적화\"", 20),
    ]:
        p = tf.add_paragraph()
        p.text = line
        p.alignment = PP_ALIGN.CENTER
        p.runs[0].font.size = Pt(size)
        p.runs[0].font.color.rgb = C_TITLE
        if size >= 20:
            p.runs[0].font.bold = True

    # 하단 시사점
    add_bullets(
        s, Inches(0.7), Inches(4.0), Inches(12), Inches(3.0),
        [
            "실무 교훈",
            ("단일 metric(δ₁ 만점) 최적화는 실제 시스템에 쓸모없을 수 있다", 1),
            ("모델 설계는 파이프라인 전체의 constraint에 종속 — 전후 모듈과 함께 설계해야", 1),
            ("'좋다고 알려진 최신 모델' 무조건 사용보다, 본인 요구사항(실시간·하드웨어·도메인)에 맞춘 조합이 유효", 1),
            ("Scaling Block처럼 '쓸 때 결정되는 파라미터'를 두면 한 번 학습한 모델의 재사용성이 폭발적으로 증가", 1),
            ("향후 본인 네트워크 설계 시: 목표 3대(속도/정확도/일반화)를 명확히 정하고, 각 구성요소 선택 근거를 논리적으로 정당화할 것", 1),
        ],
        size=13, bold_first=True,
    )

    # ========================================================================
    # 알고리즘 로직 상세 (Refinement + Avoidance)
    # ========================================================================

    # --------- 22. 섹션 구분 ------------------------------------------------
    add_section_header(
        prs,
        "알고리즘 상세 로직",
        "Depth Refinement & Obstacle Avoidance",
    )

    # --------- 23. Refinement — 왜 필요한가 + 파이프라인 -------------------
    s = add_content_slide(prs, "23. Depth Refinement — 필요성과 파이프라인")
    add_bullets(
        s, Inches(0.5), Inches(1.1), Inches(6.2), Inches(2.5),
        [
            "왜 Refinement 가 필요한가?",
            ("MDE only: 학습 분포와 다른 환경에서 스케일 오차 큼 (복도·단색 벽)", 1),
            ("Stereo only: sparse 해서 dense depth 불가, texture 약하면 불안정", 1),
            ("→ MDE 의 dense 구조 + Stereo 의 reliable metric scale 결합", 1),
        ],
        size=13, bold_first=True,
    )
    # 파이프라인 다이어그램
    tb = s.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12.5), Inches(3.5))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = "7 단계 파이프라인 (TIE 논문)"
    tf.paragraphs[0].runs[0].font.bold = True
    tf.paragraphs[0].runs[0].font.size = Pt(14)
    tf.paragraphs[0].runs[0].font.color.rgb = C_ACCENT

    for line, size, italic in [
        ("Step 1. Rectification — 좌/우 이미지를 epipolar line 수평 정렬", 12, False),
        ("Step 2. Stereo Matching (SGBM/Libelas) — 겹치는 영역에서 sparse depth 획득", 12, False),
        ("Step 3. Superpixel 분할 (SEEDS) — 시각 유사 픽셀 그룹화", 12, False),
        ("Step 4. Outlier 제거 (2 단계):", 12, False),
        ("   ① 그룹 내 stereo depth 변동 큰 그룹 제외", 11, True),
        ("   ② MDE-Stereo 중앙값이 50% 이상 차이나면 제외", 11, True),
        ("Step 5. Linear Regression (OLS): n ≥ threshold 일 때 (s_d, t_d) 추정", 12, False),
        ("Step 6. Temporal Smoothing (1차 필터): s_f ← α·s_f_prev + (1-α)·s_d", 12, False),
        ("Step 7. 전체 적용: d_refined = s_f · d̂ + t_f", 12, False),
    ]:
        p = tf.add_paragraph()
        p.text = line
        p.runs[0].font.size = Pt(size)
        p.runs[0].font.color.rgb = C_TEXT
        if italic:
            p.runs[0].font.italic = True

    # --------- 24. Refinement 수식 & Thesis vs TIE ------------------------
    s = add_content_slide(prs, "24. Refinement 수식 — Thesis (비율 스케일) vs TIE (Linear Regression)")
    # 좌측: Thesis
    tb = s.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(6.3), Inches(3.0))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = "Thesis (2022) — 비율 스케일링 (eq. 2-3)"
    tf.paragraphs[0].runs[0].font.bold = True
    tf.paragraphs[0].runs[0].font.size = Pt(14)
    tf.paragraphs[0].runs[0].font.color.rgb = C_TITLE

    for line, size, mono in [
        ("", 8, False),
        ("d^g_ref,i = r^g · (Σ d^g_stereo,c / Σ d^g_mono,c) · d^g_mono,i", 12, True),
        ("", 6, False),
        ("r^g = 1  if  d_max / d_min < 1.5", 11, True),
        ("      0  otherwise", 11, True),
        ("", 6, False),
        ("→ binary 플래그, scale 만 보정", 11, False),
        ("→ 그룹 내 stereo 분산 크면 refinement 포기", 11, False),
    ]:
        p = tf.add_paragraph()
        p.text = line if line else " "
        if line and p.runs:
            p.runs[0].font.size = Pt(size)
            if mono:
                p.runs[0].font.name = "Courier New"
            p.runs[0].font.color.rgb = C_TEXT

    # 우측: TIE
    tb = s.shapes.add_textbox(Inches(6.9), Inches(1.1), Inches(6.2), Inches(3.0))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = "TIE (2025) — 선형 회귀 (eq. 3 변형)"
    tf.paragraphs[0].runs[0].font.bold = True
    tf.paragraphs[0].runs[0].font.size = Pt(14)
    tf.paragraphs[0].runs[0].font.color.rgb = C_TITLE

    for line, size, mono in [
        ("", 8, False),
        ("d_i = s_f · d̂_i + t_f", 12, True),
        ("", 6, False),
        ("s_d, t_d = OLS regression from n inliers", 11, True),
        ("s_f ← α · s_f_prev + (1-α) · s_d", 11, True),
        ("t_f ← α · t_f_prev + (1-α) · t_d", 11, True),
        ("", 6, False),
        ("→ scale + offset 동시 보정", 11, False),
        ("→ temporal smoothing (flicker 제거)", 11, False),
    ]:
        p = tf.add_paragraph()
        p.text = line if line else " "
        if line and p.runs:
            p.runs[0].font.size = Pt(size)
            if mono:
                p.runs[0].font.name = "Courier New"
            p.runs[0].font.color.rgb = C_TEXT

    # 하단 비교표
    add_table(
        s, Inches(0.5), Inches(4.2), Inches(12.4), Inches(2.8),
        [
            ["항목", "Thesis (2022)", "TIE (2025)"],
            ["보정 모델", "d = k · d̂  (scale only)", "d = s·d̂ + t  (scale + bias)"],
            ["Outlier 처리", "binary flag r^g ∈ {0,1}", "2단계 필터링 + 조건부 업데이트"],
            ["시간적 smoothing", "없음 (flicker 발생 가능)", "1차 필터 (α smoothing)"],
            ["표현력", "scale 만 보정", "scale + offset 동시 보정"],
            ["업데이트 조건", "매 프레임 독립", "n ≥ threshold 일 때만"],
        ],
        font_size=11, header_size=12,
    )

    # --------- 25. Refinement 효과 (unseen 도메인 압도) -----------------
    s = add_content_slide(prs, "25. Refinement 효과 — Unseen Domain 에서 결정적")
    add_bullets(
        s, Inches(0.5), Inches(1.1), Inches(12.5), Inches(0.8),
        [
            "MDE only vs MDE + Refinement (TIE 논문 Table II, δ₁ 기준)",
        ],
        size=13, bold_first=True,
    )
    add_table(
        s, Inches(0.5), Inches(1.8), Inches(12.4), Inches(3.6),
        [
            ["Dataset", "환경", "MDE only δ₁", "MDE + Refinement δ₁", "향상"],
            ["KITTI (seen)",       "학습 분포",   "0.959", "0.973", "+1.5%"],
            ["vKITTI2 (unseen)",   "synthetic",   "0.860", "0.910", "+5.8%"],
            ["ApolloScape (unseen)","실외 중국",   "0.079", "0.797", "+900%"],
            ["DDAD (unseen)",      "실외 US",     "0.790", "0.843", "+6.7%"],
            ["Outdoor ≤15m (자체)", "근거리",     "0.390", "0.645", "+65%"],
        ],
        font_size=11, header_size=12,
    )
    add_bullets(
        s, Inches(0.5), Inches(5.7), Inches(12.5), Inches(1.5),
        [
            "핵심 관찰",
            ("Seen 환경: refinement 효과 작음 (이미 학습으로 잘 맞춤)", 1),
            ("Unseen 환경: refinement 가 생존 전략 — ApolloScape 에서 10배 향상", 1),
            ("→ stereo 는 '학습 없이도 동작하는 물리적 레퍼런스' 역할", 1),
            ("→ 실무에서 unseen 도메인 대응이 핵심인 로봇/UAV 에 필수", 1),
        ],
        size=13, bold_first=True,
    )

    # --------- 26. 섹션 구분: 장애물 회피 -------------------------------
    add_section_header(
        prs,
        "장애물 회피 알고리즘",
        "Behavior Arbitration + Target-point Navigation",
    )

    # --------- 27. Behavior Arbitration + Steering 수식 ------------------
    s = add_content_slide(prs, "27. Steering / Thrust 명령 생성 (TIE eq. 4-6)")
    add_bullets(
        s, Inches(0.5), Inches(1.1), Inches(12.5), Inches(1.3),
        [
            "철학: 복잡한 mapping / VIO / trajectory optimization 없이 refined depth 만으로 명령 생성",
            ("각 장애물이 밀어내는 '힘(potential)' 합성 → 방향 결정 (Althaus & Christensen 2002)", 1),
        ],
        size=12, bold_first=True,
    )
    # 수식 박스
    tb = s.shapes.add_textbox(Inches(0.5), Inches(2.6), Inches(12.5), Inches(4.0))
    tf = tb.text_frame
    tf.word_wrap = True
    for line, size, bold, color_override, mono in [
        ("Exponential depth (가까울수록 큰 가중치)", 13, True, C_ACCENT, False),
        ("   D_g = exp(d_c − d_g) · N_g · S_n / (H · W)", 12, False, C_TEXT, True),
        ("      d_c : user-defined 기준 거리", 10, False, C_SUB, False),
        ("      N_g : 그룹 픽셀 수 (큰 장애물 = 큰 가중)", 10, False, C_SUB, False),
        ("      S_n : 전체 superpixel 개수 (정규화)", 10, False, C_SUB, False),
        ("", 6, False, C_TEXT, False),
        ("Horizontal / Vertical 집계 (이미지 중심 강조)", 13, True, C_ACCENT, False),
        ("   φ_f^h = Σ D_g · exp(-(u_g / HFOV)²)", 12, False, C_TEXT, True),
        ("   φ_f^v = Σ D_g · exp(-(v_g / VFOV)²)", 12, False, C_TEXT, True),
        ("", 6, False, C_TEXT, False),
        ("Smoothing (명령 폭주 방지)", 13, True, C_ACCENT, False),
        ("   δ_h = φ_f^h · exp(-(φ_f^h / G_a)²)       ← steering", 12, False, C_TEXT, True),
        ("   δ_v = φ_f^v · exp(-(φ_f^v / G_a)²)       ← altitude rate", 12, False, C_TEXT, True),
    ]:
        if tf.text == "" and not tf.paragraphs[0].text:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line if line else " "
        if line and p.runs:
            p.runs[0].font.size = Pt(size)
            p.runs[0].font.color.rgb = color_override
            p.runs[0].font.bold = bold
            if mono:
                p.runs[0].font.name = "Courier New"

    # --------- 28. Algorithm 1 pseudocode -------------------------------
    s = add_content_slide(prs, "28. Algorithm 1 — Target-point Navigation (TIE 논문)")
    # pseudocode 박스
    tb = s.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12.5), Inches(5.8))
    tf = tb.text_frame
    tf.word_wrap = True
    for line, size, bold, color in [
        ("목적: 회피 + 목표점 주행을 동시에 수행하는 자율 비행", 13, True, C_ACCENT),
        ("", 6, False, C_TEXT),
        ("function get_steering_and_collision_prob(superpixels):", 11, True, C_TITLE),
        ("    δ_h, δ_v = eq.(6)  # smoothed 회피 명령", 10, False, C_TEXT),
        ("    δ_coll = (φ_f^h + φ_f^v) / S_n  # 전체 장면 위험도", 10, False, C_TEXT),
        ("    p_coll = δ_coll · exp((δ_coll · G_c)²)  # 충돌 확률 ∈ [0,1]", 10, False, C_TEXT),
        ("    return δ_h, δ_v, p_coll", 10, False, C_TEXT),
        ("", 6, False, C_TEXT),
        ("function command_for_desired_location(δ_hl, δ_hr, p_coll_l, p_coll_r, P_goal, P_cur):", 11, True, C_TITLE),
        ("    p_coll_t = 0.3 · mean(p_coll_l, p_coll_r) + 0.7 · p_coll_t_prev", 10, False, C_TEXT),
        ("    v_d = V_max · (1 − p_coll_t)  # 충돌 위험 ↑ → 속도 ↓", 10, False, C_TEXT),
        ("    if distance(P_cur, P_goal) > threshold:", 10, False, C_TEXT),
        ("        δ_g = desired_heading(P_cur, P_goal)  # 목표점 방향", 10, False, C_TEXT),
        ("        δ_t = δ_hr − δ_hl  # 좌/우 회피 차이", 10, False, C_TEXT),
        ("        w, tc = get_weight(δ_t)  # 회피/목표 블렌딩 비율", 10, False, C_TEXT),
        ("        w_time = exp(-0.05·(t_now - tc)) · w  # 시간 감쇠", 10, False, C_TEXT),
        ("        δ_blend = w_time · δ_t + (1 − w_time) · δ_g", 10, False, C_TEXT),
        ("        ψ_t = 0.7 · ψ_t_prev + 0.3 · δ_blend  # heading low-pass", 10, False, C_TEXT),
        ("    return ψ_t, v_d", 10, False, C_TEXT),
        ("", 6, False, C_TEXT),
        ("function get_weight(δ_t):", 11, True, C_TITLE),
        ("    if |δ_t| < 10°:  return w=0.5, tc: no update  # 작은 장애물 → 반반", 10, False, C_TEXT),
        ("    else: return w = min(|δ_t|/20, 1), tc = t_now  # 큰 장애물 → 회피 우선", 10, False, C_TEXT),
    ]:
        if tf.text == "" and not tf.paragraphs[0].text:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line if line else " "
        if line and p.runs:
            p.runs[0].font.size = Pt(size)
            p.runs[0].font.color.rgb = color
            p.runs[0].font.bold = bold
            p.runs[0].font.name = "Courier New" if size == 10 else "Calibri"

    # ========================================================================
    # 실험 결과
    # ========================================================================

    # --------- 29. 섹션 구분 --------------------------------------------
    add_section_header(
        prs,
        "본 프로젝트 실험 결과",
        "KITTI · NYU 학습 및 정량/정성 평가",
    )

    # --------- 30. KITTI 학습 곡선 --------------------------------------
    s = add_content_slide(prs, "30. KITTI 학습 결과 (7 Epochs, RTX 5070)")
    add_bullets(
        s, Inches(0.5), Inches(1.1), Inches(6.0), Inches(1.2),
        [
            "학습 설정",
            ("batch=8, AdamW lr=1e-4, cosine schedule, SI loss (α=10, λ=0.85)", 1),
            ("총 44,008 train / 1,734 val / 652 test (Eigen split)", 1),
        ],
        size=12, bold_first=True,
    )
    # 좌측 loss 추이
    add_table(
        s, Inches(0.3), Inches(2.6), Inches(6.0), Inches(4.2),
        [
            ["Epoch", "Train Loss", "Val Loss", "LR"],
            ["1", "1.176", "0.787", "9.5e-5"],
            ["2", "0.746", "0.648", "8.1e-5"],
            ["3", "0.615", "0.583", "6.1e-5"],
            ["4", "0.534", "0.497", "3.9e-5"],
            ["5", "0.476", "0.450", "1.9e-5"],
            ["6", "0.437", "0.427", "5e-6"],
            ["7", "0.417", "0.413", "0"],
        ],
        font_size=11, header_size=12,
    )
    # 우측 관찰
    add_bullets(
        s, Inches(6.6), Inches(2.6), Inches(6.4), Inches(4.4),
        [
            "관찰",
            ("Loss 꾸준히 감소 (1.18 → 0.42) — underfitting 없음", 1),
            ("Train/Val gap 거의 없음 — overfitting 없음", 1),
            ("Cosine scheduler 정상 동작 (lr 0 까지 감소)", 1),
            ("", 1),
            "속도 (RTX 5070)",
            ("2.17 iter/s, epoch 당 42.5 분", 1),
            ("25 epoch 완주 시 약 18 h", 1),
            ("7 epoch → 5 h (효율적 분석 가능)", 1),
        ],
        size=13, bold_first=True,
    )

    # --------- 31. KITTI Test Set 평가 + 벤치마크 -----------------------
    s = add_content_slide(prs, "31. KITTI Test Set 평가 — 논문 대비 경쟁력")
    add_table(
        s, Inches(0.3), Inches(1.1), Inches(12.7), Inches(3.0),
        [
            ["메트릭", "본 프로젝트 (7ep)", "TIE 논문 (25ep)", "해석"],
            ["δ₁ (< 1.25)",   "0.9589", "0.959", "95.9% 픽셀 정확 범위 — 논문 수준"],
            ["δ₂ (< 1.25²)",  "0.9940", "—",     "99.4%"],
            ["δ₃ (< 1.25³)",  "0.9986", "—",     "99.9%"],
            ["AbsRel",        "0.0614", "0.065", "상대 오차 6.1% — 논문 약간 상회"],
            ["RMSE",          "3.01 m", "2.44 m","평균 3m 오차 (25ep 시 개선 여지)"],
        ],
        font_size=12, header_size=13,
    )

    # KITTI 리더보드 비교
    add_bullets(
        s, Inches(0.5), Inches(4.3), Inches(12.5), Inches(0.5),
        [
            "KITTI Eigen split 리더보드 포지션 (δ₁ 기준)",
        ],
        size=13, bold_first=True,
    )
    add_table(
        s, Inches(0.3), Inches(4.8), Inches(12.7), Inches(2.5),
        [
            ["Rank", "Method", "Year", "δ₁", "파라미터"],
            ["1",  "SPIdepth",            "2024", "0.990", "—"],
            ["2",  "UniK3D",              "2025", "0.990", "—"],
            ["3",  "Metric3Dv2 (ViT-g2)", "2024", "0.989", "ViT-g2 (3B+)"],
            ["-",  "Depth Anything V2-L", "2024", "0.982", "335M"],
            ["-",  "NeWCRFs",             "2022", "0.974", "270M"],
            ["-",  "AdaBins",             "2021", "0.964", "78M"],
            ["-",  "Ours (ConvNeXt+LWA)", "2025", "0.959", "15M"],
            ["-",  "BTS",                 "2019", "0.956", "47M"],
        ],
        font_size=10, header_size=11,
    )

    # --------- 32. NYU 학습 곡선 -----------------------------------------
    s = add_content_slide(prs, "32. NYU 학습 결과 (7 Epochs, RTX 5070)")
    add_bullets(
        s, Inches(0.5), Inches(1.1), Inches(6.0), Inches(1.2),
        [
            "학습 설정",
            ("batch=8, AdamW lr=1e-4, cosine, max_depth=10m", 1),
            ("47,584 train / 654 val (FastDepth 전처리판)", 1),
        ],
        size=12, bold_first=True,
    )
    add_table(
        s, Inches(0.3), Inches(2.6), Inches(6.0), Inches(4.2),
        [
            ["Epoch", "Train Loss", "Val Loss"],
            ["1",  "1.449", "1.702"],
            ["2",  "0.920", "1.654"],
            ["3",  "0.766", "1.585"],
            ["4",  "0.669", "1.578"],
            ["5",  "0.603", "1.589"],
            ["6",  "0.555", "1.576 ← 최적"],
            ["7",  "0.531", "1.597 (overfit 시작)"],
        ],
        font_size=11, header_size=12,
    )
    add_bullets(
        s, Inches(6.6), Inches(2.6), Inches(6.4), Inches(4.4),
        [
            "관찰",
            ("Epoch 6 에서 val loss 최저 (1.576)", 1),
            ("Epoch 7 에서 약간 상승 — mild overfitting 시작", 1),
            ("Train/Val gap 큼 (실내 장면 다양성 ↑)", 1),
            ("", 1),
            "KITTI 대비 차이점",
            ("NYU val loss 가 KITTI 보다 훨씬 높음", 1),
            ("실내 texture 다양성 >> 실외 도로 반복 패턴", 1),
            ("추가 augmentation / 더 많은 epoch 필요할 수 있음", 1),
        ],
        size=13, bold_first=True,
    )

    # --------- 33. NYU 평가 + domain matching 시각 비교 -----------------
    s = add_content_slide(prs, "33. NYU Val Set 평가 + Domain Matching 증명")
    add_table(
        s, Inches(0.3), Inches(1.1), Inches(12.7), Inches(3.0),
        [
            ["메트릭", "본 프로젝트 (7ep, ep6)", "AdaBins (25ep)", "BTS", "DenseDepth"],
            ["δ₁",     "0.8592", "0.903", "0.885", "0.846"],
            ["δ₂",     "0.9750", "—",     "—",     "—"],
            ["δ₃",     "0.9938", "—",     "—",     "—"],
            ["AbsRel", "0.1203", "0.103", "0.110", "0.123"],
            ["RMSE",   "0.4549 m","—",    "—",     "—"],
        ],
        font_size=11, header_size=12,
    )
    add_bullets(
        s, Inches(0.5), Inches(4.3), Inches(12.5), Inches(2.5),
        [
            "실내 이미지 추론 시 Domain Matching 의 중요성",
            ("동일한 네트워크 구조, 다른 학습 데이터셋 → 극명한 차이", 1),
            ("  · KITTI-trained (max 80m) 로 실내 이미지 추론: 3.26 ~ 80.0 m (일부 포화, 도로로 오해)", 1),
            ("  · NYU-trained   (max 10m) 로 실내 이미지 추론: 1.43 ~ 3.48 m (정확한 실내 스케일)", 1),
            ("결론: 네트워크 구조보다 '학습 데이터의 domain 일치' 가 결정적", 1),
            ("→ 향후 Phase 3 refinement 가 이 domain gap 을 해소할 것", 1),
        ],
        size=13, bold_first=True,
    )

    # --------- 34. Ubuntu 원격 제어 워크플로우 --------------------------
    s = add_content_slide(prs, "34. Ubuntu 원격 제어 워크플로우 — 학습 자동화")
    # 다이어그램 영역
    tb = s.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12.5), Inches(2.5))
    tf = tb.text_frame
    tf.word_wrap = True
    for line, size, bold, mono in [
        ("Mac Mini (개발)  ─────  GitHub  ─────  Ubuntu (RTX 5070, 학습)", 13, True, False),
        ("     ①코드 작성     ②push/pull        ③SSH 명령 + nohup", 11, False, False),
        ("     ⑥ scp          ──────  weight 파일 복사  ──────  ④학습 진행", 11, False, False),
        ("                                                        ⑤학습 완료", 11, False, False),
    ]:
        if tf.text == "" and not tf.paragraphs[0].text:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line if line else " "
        if line and p.runs:
            p.runs[0].font.size = Pt(size)
            p.runs[0].font.color.rgb = C_TEXT
            p.runs[0].font.bold = bold
            if mono:
                p.runs[0].font.name = "Courier New"
            p.alignment = PP_ALIGN.CENTER

    add_bullets(
        s, Inches(0.5), Inches(3.8), Inches(12.5), Inches(3.5),
        [
            "주요 기법",
            ("SSH Key 인증 (ssh-keygen + ssh-copy-id) — 비밀번호 없이 자동화", 1),
            ("nohup + '&' — SSH 끊겨도 수 시간 학습 계속 진행", 1),
            ("로그 리다이렉트 '> log 2>&1' — stdout/stderr 기록, tail 로 실시간 모니터링", 1),
            ("GitHub 경유 코드 동기화 — Mac 에서 push, Ubuntu 에서 pull (버전 추적)", 1),
            ("scp — 대용량 weight (100MB+) 파일을 Ubuntu → Mac 복사", 1),
            "이점",
            ("Mac Mini 꺼도 학습 계속 (퇴근 후 밤새 학습)", 1),
            ("대용량 데이터 (KITTI 175GB, NYU 35GB) 는 Ubuntu 가 직접 다운로드 — Mac 거치지 않음", 1),
            ("세션 기반 명령 실행 → Context 낭비 없음", 1),
        ],
        size=13, bold_first=True,
    )

    # --------- 35. 섹션 구분: 결론 -------------------------------------
    add_section_header(
        prs,
        "결론 및 향후 방향",
        "실무 가이드 · 트렌드 · 인사이트",
    )

    # --------- 36. 실무 선택 가이드 ----------------------------------------
    s = add_content_slide(prs, "36. 실무 선택 가이드")
    add_table(
        s, Inches(0.4), Inches(1.2), Inches(12.5), Inches(4.5),
        [
            ["시나리오", "추천 모델", "이유"],
            ["UAV / 모바일 / 임베디드",
             "ConvNeXt+LWA · Depth Anything-S · Monodepth2",
             "15–25M params, 10–50ms inference"],
            ["실시간 로봇 perception",
             "ConvNeXt+LWA · GLPDepth · BTS",
             "실시간 + 고정확도 균형"],
            ["KITTI·NYU SOTA 추구",
             "Depth Anything v2-L · NewCRFs · ZoeDepth",
             "정확도 최상위, GPU 필요"],
            ["Unseen domain (일반화)",
             "Depth Anything v2 · ZoeDepth · MiDaS",
             "다양한 데이터 pretrain"],
            ["GT 없이 학습",
             "Monodepth2 · SfMLearner",
             "photometric loss 기반 self-sup."],
            ["최고 품질, offline",
             "Marigold · Depth Pro · DepthFM",
             "diffusion 기반 sharp edge"],
            ["Rapid prototyping",
             "HuggingFace Depth Anything v2",
             "one-line pipeline (transformers)"],
        ],
        font_size=12, header_size=13,
    )

    # --------- 23. 향후 트렌드 --------------------------------------------
    s = add_content_slide(prs, "37. 향후 트렌드 (2025+)")
    add_bullets(
        s, Inches(0.6), Inches(1.2), Inches(12), Inches(5.8),
        [
            "Foundation → Task-specific fine-tuning 패러다임 정착",
            ("DINOv2 / SAM / Depth Anything backbone + 얇은 task head", 1),
            ("적은 labeled 데이터로도 도메인 특화 가능", 1),
            "Video / Temporal consistency",
            ("Video Depth Anything, Online VDA (2025) 등 long-term consistency 연구", 1),
            ("프레임 간 flicker 제거, tracking/SLAM 친화", 1),
            "Metric scale recovery",
            ("relative depth foundation + metric scale alignment 방식이 주류", 1),
            ("ZoeDepth → UniDepth → Depth Pro (metric zero-shot)", 1),
            "Flow matching & 빠른 diffusion",
            ("Marigold-LCM, DepthFM: diffusion 품질 + real-time 속도", 1),
            "센서 융합",
            ("sparse LiDAR + RGB = depth completion (계속 연구 중)", 1),
            ("event camera + RGB (Dynamic obstacle)", 1),
            "경량화 & on-device",
            ("MobileViT, EfficientFormer + quantization → smartphone 실시간", 1),
        ],
        size=15, bold_first=True,
    )

    # --------- 15. 인사이트 요약 -------------------------------------------
    s = add_content_slide(prs, "38. 주요 인사이트 요약")
    add_bullets(
        s, Inches(0.6), Inches(1.2), Inches(12), Inches(5.8),
        [
            "1. 더 크면 더 정확, 하지만 Pareto 최적은 중·경량대",
            ("15M (ConvNeXt+LWA) vs 335M (Depth Anything-L) 정확도 차이 < 3%", 1),
            ("실제 운용에선 경량 모델이 유리한 경우가 많음", 1),
            "2. Encoder 선택이 절대적으로 중요",
            ("Decoder/Head는 정제에 가깝고, feature 품질은 encoder가 결정", 1),
            ("Foundation encoder(DINOv2) + 얇은 head = 가장 강력한 조합", 1),
            "3. Zero-shot 일반화가 새 전쟁터",
            ("single-dataset SOTA의 시대는 끝, multi-dataset robustness 중심", 1),
            ("unseen domain에서 무너지면 서비스 불가", 1),
            "4. Loss는 크게 바뀌지 않음",
            ("SILog + gradient matching이 여전히 표준", 1),
            "5. 실시간 UAV/robotics에서는 ConvNeXt+LWA 류 경량 구조 유망",
            ("실험에서 7 epoch만에 논문급 정확도 달성 (δ₁=0.959) → 파이프라인 견고", 1),
            ("추가 학습·Depth Refinement로 실내 일반화도 가능할 전망", 1),
        ],
        size=14, bold_first=True,
    )

    # --------- 16. 참고 문헌 -----------------------------------------------
    s = add_content_slide(prs, "39. 참고 문헌 및 리소스")
    add_bullets(
        s, Inches(0.6), Inches(1.2), Inches(12.2), Inches(5.8),
        [
            "주요 논문",
            ("Eigen et al., Depth Map Prediction from a Single Image, NIPS 2014", 1),
            ("Laina et al., FCRN, 3DV 2016", 1),
            ("Fu et al., DORN, CVPR 2018", 1),
            ("Alhashim & Wonka, DenseDepth, arXiv 2018", 1),
            ("Godard et al., Monodepth2, ICCV 2019", 1),
            ("Lee et al., BTS, arXiv 2019", 1),
            ("Bhat et al., AdaBins, CVPR 2021", 1),
            ("Ranftl et al., DPT, ICCV 2021", 1),
            ("Kim et al., GLPDepth, arXiv 2022", 1),
            ("Yuan et al., NeWCRFs, CVPR 2022", 1),
            ("Bhat et al., ZoeDepth, arXiv 2023", 1),
            ("Yang et al., Depth Anything v1/v2, CVPR / NeurIPS 2024", 1),
            ("Ke et al., Marigold, CVPR 2024 (Oral)", 1),
            "",
            "본 프로젝트의 두 논문 (원저자: Hyeongjin Kim)",
            ("[Thesis] H. Kim, \"Refined Depth Estimation and Safety Navigation with a Binocular Camera\", M.S. Thesis, KNU, 2022", 1),
            ("[TIE] E. Cho, H. Kim, P. Kim, H. Lee, \"Obstacle Avoidance of a UAV Using Fast Monocular Depth Estimation for a Wide Stereo Camera\", IEEE TIE 72(2), 2025", 1),
            "",
            "코드 & 벤치마크",
            ("KITTI Depth Prediction: https://www.cvlibs.net/datasets/kitti/eval_depth.php", 1),
            ("Monocular Depth Estimation Toolbox (zhyever): 통합 벤치마크", 1),
            ("timm (Hugging Face): ConvNeXt v2 등 최신 backbone 라이브러리", 1),
        ],
        size=12, bold_first=True,
    )

    return prs


if __name__ == "__main__":
    import os
    os.makedirs("docs/mde_survey", exist_ok=True)
    prs = build()
    out = "docs/mde_survey/MDE_Network_Survey.pptx"
    prs.save(out)
    print(f"Saved: {out}")
    print(f"Slides: {len(prs.slides)}")
